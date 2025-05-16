import json
from pptx import Presentation
import argparse

import calcslidesimilarity

# pptx のスライドの幅と高さを取得するための変数
slide_width = 0
slide_height = 0

#---------------------------------------------
# 当スクリプト動作設定用パラメータ
#---------------------------------------------
# スライド類似度閾値。この数値未満なら類似とはみなさない
slidesimilarity_threshold = 0.7

#---------------------------------------------
# calcslidesimilarity モジュールの動作設定用パラメータ
#---------------------------------------------
# テキストの完全一致を要求する場合は True にする
text_strictmatch = True
# レイアウトとテキストの類似度を計算するための重み
# 合計で1.0 になるように設定
layout_weight = 0.5
text_weight = 0.5
# シェイプの類似度の閾値
shape_threshold = 0.8  # 1.0なら完全一致していなければ類似とみなさない
# スライドの類似度の閾値
slide_threshold = 0.8
############### 注意！！！ ###############
# slide_threshold は find_similar_slide_pairs() で使用するが、
# 当スクリプトでは今のところ 同関数を使用していないので、
# slide_threshold は使用されていない。
#########################################

# calcslidesimilarity モジュールに動作設定用パラメータをセットする
calcslidesimilarity.set_similarity_settings(
    layout=layout_weight,
    text=text_weight,
    slide=slide_threshold,
    shape=shape_threshold,
    text_strict=text_strictmatch)


#---------------------------------------------
# コマンドライン引数を解析して返す
#---------------------------------------------
def parse_commandargs():
    # コマンドライン引数のパーサーを作成
    parser = argparse.ArgumentParser(description="Analyze PowerPoint file and output it as .json file.")
    # ----------------------------------------------
    # yaml指定時とpptx単独指定時に共通する引数の定義
    parser.add_argument("file1", help="PPTX file to be analyzed.")
    parser.add_argument("file2", help="PPTX file to be analyzed.")
    # ログレベル指定　オプション --loglevel ,文字型、デフォルトは info
    parser.add_argument('--loglevel',  type=str, default='info', help='log level [debug|INFO]')
    # ログファイル指定　オプション --lf, 略称 -f,文字型、デフォルトは STDOUT
    parser.add_argument('--logfile',  type=str, default='STDOUT', help='path to log file, or STDOUT if omitted')
    # ----------------------------------------------
    # ソースとエクスポートのフォルダー、ファイル指定
    parser.add_argument('--exportdir', '-e', type=str, default='', help='export directory for index file')
    parser.add_argument('--exportfile', '-o', type=str, default=None, help='export file name for index file')
    parser.add_argument('--sourcedir', '-s', type=str, default='', help='source directory for source file')


    # 引数を解析
    _args = parser.parse_args()
    return _args




#---------------------------------------------
# shape の座標値を辞書形式で取得する
#---------------------------------------------
def shape_position_dict(shape):
    return {
        "left": shape.left.pt,
        "top": shape.top.pt,
        "width": shape.width.pt,
        "height": shape.height.pt
    }

#---------------------------------------------
#座標値をスライドの大きさとの相対比率に変換する
#---------------------------------------------
def normalize_position(shape, slide_width_pt, slide_height_pt):
    return {
        "left": round(shape.left.pt / slide_width_pt, 4),
        "top": round(shape.top.pt / slide_height_pt, 4),
        "width": round(shape.width.pt / slide_width_pt, 4),
        "height": round(shape.height.pt / slide_height_pt, 4),
    }

#---------------------------------------------
# shape["position_ratio"] 以下の座標値を上の階層に移動して
# position_ratio と position_pt を削除する
#---------------------------------------------
def position_ratio_to_upstair(shape):
    # shape.left.pt, shape.top.pt, shape.width.pt, shape.height.pt
    # から left, top, width, height を計算する
    shape["left"] = shape["position_ratio"]["left"]
    shape["top"] = shape["position_ratio"]["top"]
    shape["width"] = shape["position_ratio"]["width"]
    shape["height"] = shape["position_ratio"]["height"]

    # shapeから ["position_ratio"],["position_pt"] を削除
    del shape["position_ratio"]
    del shape["position_pt"]

    return shape

#---------------------------------------------
# pptx ファイルを解析して、スライドの情報を取得する
# 解析結果を辞書形式で返す
#---------------------------------------------
def analyze_pptx(filepath):
    prs = Presentation(filepath)

    slide_width = prs.slide_width.pt
    slide_height = prs.slide_height.pt


    # 表紙タイトル（slide 0）
    document_title = ""
    if prs.slides:
        first_slide = prs.slides[0]
        max_font_size = 0
        for shape in first_slide.shapes:
            if not hasattr(shape, "text") or not shape.text.strip():
                continue
            # テキストがある shape を対象
            try:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_size = run.font.size.pt if run.font.size else 0
                        if font_size > max_font_size:
                            max_font_size = font_size
                            document_title = run.text.strip()
            except AttributeError:
                continue  # shape に text_frame がない場合もある

    analyzed = {
        "DocumentTitle": document_title,
        "slides": []
    }

    for idx, slide in enumerate(prs.slides):
        slide_title = ""
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:
                slide_title = shape.text.strip()
                break

        shapes = []
        for shape in slide.shapes:
            # pos = shape_position_dict(shape)

            pos_pt = shape_position_dict(shape)
            pos_ratio = normalize_position(shape, slide_width, slide_height)


            # テキスト
            if hasattr(shape, "text") and shape.text.strip():
                shapes.append({
                    "type": "text",
                    "text": shape.text.strip(),
                    "position_pt": pos_pt,
                    "position_ratio": pos_ratio
                })

            # 画像
            elif shape.shape_type == 13 and shape.image is not None:
                image_bytes = shape.image.blob
                ext = shape.image.ext
                image_name = shape.image.filename  # e.g. image1.png
                shapes.append({
                    "type": "image",
                    "image_format": ext,
                    "filename": image_name,
                    "image_bytes": len(image_bytes),
                    "position_pt": pos_pt,
                    "position_ratio": pos_ratio
                })

            # 表
            elif shape.has_table:
                table_data = []
                table = shape.table
                for row in table.rows:
                    table_data.append([cell.text.strip() for cell in row.cells])
                shapes.append({
                    "type": "table",
                    "rows": table_data,
                    "position_pt": pos_pt,
                    "position_ratio": pos_ratio
                })

        # ノート
        notes_text = ""
        if slide.has_notes_slide and hasattr(slide.notes_slide, "notes_text_frame"):
            frame = slide.notes_slide.notes_text_frame
            if frame is not None:
                notes_text = frame.text.strip()

        analyzed["slides"].append({
            "slidetitle": slide_title,
            "shapes": shapes,
            "notes": notes_text
        })

    return analyzed

# 使用例
if __name__ == "__main__":

    args = parse_commandargs()

    # コマンドライン引数の取得
    exportdir = args.exportdir
    exportfile = args.exportfile
    sourcedir = args.sourcedir
    pptxfile1 = args.file1
    pptxfile2 = args.file2
    logfile = args.logfile
    loglevel = args.loglevel

    pptxmap1 = analyze_pptx(pptxfile1)
    # print(json.dumps(pptxmap1, ensure_ascii=False, indent=2))

    #pptxmap1["slides"][x]["shapes"] を position_ratio_to_upstair で変換して更新する
    slides1 = pptxmap1["slides"]
    for slide in slides1:
        for shape in slide["shapes"]:
            shape = position_ratio_to_upstair(shape)

    pptxmap2 = analyze_pptx(pptxfile2)
    # print(json.dumps(pptxmap2, ensure_ascii=False, indent=2))
    slides2 = pptxmap2["slides"]
    for slide in slides2:
        for shape in slide["shapes"]:
            shape = position_ratio_to_upstair(shape)


    # pptx1 のスライドの類似データを格納するリスト
    sm1item = {}
    for  slide1idx in range( 1 , len(slides1)):
        sm1item[f"{slide1idx}"] = {}

    # pptx2 のスライドの類似データを格納するリスト
    sm2item = {}
    for  slide2idx in range( 1 , len(slides2)):
        sm2item[f"{slide2idx}"] = {}

    for  slide1idx in range( 1 , len(slides1)):
        slide1 = pptxmap1["slides"][slide1idx]
        for slide2idx in range(1, len(slides2)):
            slide2 = pptxmap2["slides"][slide2idx]
            # スライドの類似度を計算
            similarity = calcslidesimilarity.slide_similarity(slide1, slide2, shape_threshold)
            if( similarity >= slidesimilarity_threshold):
                sm1item[f"{slide1idx}"][f"{slide2idx}"]=similarity
                sm2item[f"{slide2idx}"][f"{slide1idx}"]=similarity

    # 類似度が高い順にソートする
    for slide1idx in range( 1 , len(slides1)):
        sm1item[f"{slide1idx}"] = sorted(sm1item[f"{slide1idx}"].items(), key=lambda x: x[1], reverse=True)
    for slide2idx in range( 1 , len(slides2)):
        sm2item[f"{slide2idx}"] = sorted(sm2item[f"{slide2idx}"].items(), key=lambda x: x[1], reverse=True)

    #print(json.dumps(sm1item, ensure_ascii=False, indent=2))

    print(f"==========  一致度 = 1.0 のスライドを表示します ==========")
    for slide1idx in range( 1 , len(slides1)):
        for z in range(len(sm1item[f"{slide1idx}"])):
            # 類似度が1のものを出力する
            # print( f"類似度 {slide1idx} {z} {sm1item[str(slide1idx)][z][1]}")
            if sm1item[f"{slide1idx}"][z][1] == 1.0:
                # print(f"{sm1item[str(slide1idx)][z]}")
                slide2idx = int(sm1item[str(slide1idx)][z][0])
                print('-' * 20)
                print(f"{slide1idx} {slides1[slide1idx]["slidetitle"]}")
                print(f"{slide2idx} {slides2[slide2idx]["slidetitle"]}")

    print('\n')

    print(f"==========  一致度 0.8 =< ～ < 1.0 のスライドを表示します ==========")
    for slide1idx in range( 1 , len(slides1)):
        for z in range(len(sm1item[f"{slide1idx}"])):
            # 類似度が1のものを出力する
            # print( f"類似度 {slide1idx} {z} {sm1item[str(slide1idx)][z][1]}")
            if sm1item[f"{slide1idx}"][z][1] >= 0.8 and sm1item[f"{slide1idx}"][z][1] <1.0 :
                # print(f"{sm1item[str(slide1idx)][z]}")
                slide2idx = int(sm1item[str(slide1idx)][z][0])
                print('-' * 20)
                print(f"一致度 {sm1item[f'{slide1idx}'][z][1]}")
                print(f"{slide1idx} {slides1[slide1idx]["slidetitle"]}")
                print(f"{slide2idx} {slides2[slide2idx]["slidetitle"]}")

    print(f"==========  一致度 0.6 =< ～ < 0.8 のスライドを表示します ==========")
    for slide1idx in range( 1 , len(slides1)):
        for z in range(len(sm1item[f"{slide1idx}"])):
            # 類似度が1のものを出力する
            # print( f"類似度 {slide1idx} {z} {sm1item[str(slide1idx)][z][1]}")
            if sm1item[f"{slide1idx}"][z][1] >= 0.6 and sm1item[f"{slide1idx}"][z][1] <0.8 :
                # print(f"{sm1item[str(slide1idx)][z]}")
                slide2idx = int(sm1item[str(slide1idx)][z][0])
                print('-' * 20)
                print(f"一致度 {sm1item[f'{slide1idx}'][z][1]}")
                print(f"{slide1idx} {slides1[slide1idx]["slidetitle"]}")
                print(f"{slide2idx} {slides2[slide2idx]["slidetitle"]}")

    print(f"==========  類似スライド不検出のスライドを表示します ==========")
    for slide1idx in range( 1 , len(slides1)):
        # sm1item[f"{slide1idx}"] が空のものを出力する
        if len(sm1item[f"{slide1idx}"]) == 0:
            print(f"{slide1idx} {slides1[slide1idx]["slidetitle"]}")

#    print(json.dumps(sm1item, ensure_ascii=False, indent=2))
#    print(json.dumps(sm2item, ensure_ascii=False, indent=2))

