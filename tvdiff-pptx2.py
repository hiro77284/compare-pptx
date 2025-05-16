import os
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity

# モデル準備
model = SentenceTransformer('all-MiniLM-L6-v2')

# スライドサイズ（PowerPointデフォルト、ポイント単位）
SLIDE_WIDTH = 914400
SLIDE_HEIGHT = 514350

# PowerPointファイルからテキスト＋位置情報を抽出
def extract_slide_shapes(pptx_path):
    prs = Presentation(pptx_path)
    slide_vectors = []

    for slide in prs.slides:
        shape_vectors = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if not text:
                continue

            # テキストベクトル
            text_vec = model.encode(text)

            # 位置ベクトル（正規化）
            if hasattr(shape, 'left'):
                pos_vec = np.array([
                    shape.left / SLIDE_WIDTH,
                    shape.top / SLIDE_HEIGHT,
                    shape.width / SLIDE_WIDTH,
                    shape.height / SLIDE_HEIGHT
                ])
            else:
                pos_vec = np.zeros(4)

            # 結合ベクトル（テキスト + 位置）
            combined = np.concatenate([text_vec, pos_vec])
            shape_vectors.append(combined)

        slide_vectors.append(shape_vectors)
    return slide_vectors  # List[List[np.ndarray]] 各スライドのshapeベクトル群

# shape同士の類似度を比較してマッチ数を数える
def slide_similarity(vecs1, vecs2, threshold=0.85):
    if not vecs1 or not vecs2:
        return 0.0

    sim_matrix = cosine_similarity(vecs1, vecs2)

    matched = 0
    used_b_indices = set()

    for i in range(len(vecs1)):
        # 各 shape A[i] について最も近い shape B[j] を探す
        best_j = np.argmax(sim_matrix[i])
        if sim_matrix[i][best_j] > threshold and best_j not in used_b_indices:
            matched += 1
            used_b_indices.add(best_j)

    total_possible = min(len(vecs1), len(vecs2))
    return matched / total_possible if total_possible else 0.0

# ファイル間の全スライド類似度を計算
def compare_presentations(file_a, file_b):
    slides_a = extract_slide_shapes(file_a)
    slides_b = extract_slide_shapes(file_b)

    results = []
    for i, vecs_a in enumerate(slides_a):
        for j, vecs_b in enumerate(slides_b):
            score = slide_similarity(vecs_a, vecs_b)
            results.append((i+1, j+1, score))
    return results

# 使用例
if __name__ == "__main__":
    pptx_a = "example/tvdiffsample1.pptx"
    pptx_b = "example/tvdiffsample2.pptx"

    scores = compare_presentations(pptx_a, pptx_b)
    for i, j, s in scores:
        print(f"Slide A#{i} vs Slide B#{j} → 類似度: {s:.2f}")
