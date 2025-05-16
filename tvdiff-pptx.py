from pptx import Presentation
from sentence_transformers import SentenceTransformer
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import sys

# モデルの準備
model = SentenceTransformer('all-MiniLM-L6-v2')

# スライドサイズ（ポイント単位、デフォルト16:9 = 914400 x 514350）
SLIDE_WIDTH = 914400
SLIDE_HEIGHT = 514350

def slide_to_vector(pptx_path):
    prs = Presentation(pptx_path)
    slide_vectors = []

    for slide in prs.slides:
        shape_vectors = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if not text:
                continue

            # テキストの意味ベクトル
            text_vec = model.encode(text)

            # 位置ベクトル（正規化）
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                pos_vec = np.array([
                    shape.left / SLIDE_WIDTH,
                    shape.top / SLIDE_HEIGHT,
                    shape.width / SLIDE_WIDTH,
                    shape.height / SLIDE_HEIGHT
                ])
            else:
                pos_vec = np.zeros(4)

            # 結合：意味 + 位置
            shape_vec = np.concatenate([text_vec, pos_vec])
            shape_vectors.append(shape_vec)

        # スライド全体を1つのベクトルに集約（平均）
        if shape_vectors:
            slide_vector = np.mean(shape_vectors, axis=0)
            slide_vectors.append(slide_vector)

    return slide_vectors  # 複数スライドある場合もあるためリストで返す

# 使用例
vecs_a = slide_to_vector(sys.argv[1])
vecs_b = slide_to_vector(sys.argv[2])

# 最初のスライド同士を比較（例）
sim = cosine_similarity([vecs_a[0]], [vecs_b[0]])
print(f"類似度: {sim[0][0]:.3f}")

