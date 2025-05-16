from pptx import Presentation
from sentence_transformers import SentenceTransformer
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import sys

# モデルの準備
model = SentenceTransformer('all-MiniLM-L6-v2')

# PowerPointファイルからスライドごとのテキストを抽出
def extract_slide_texts(pptx_path):
    prs = Presentation(pptx_path)
    slide_texts = []
    
    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if text:
                slide_text += text + " "  # すべてのテキストを1つの文字列として結合
        slide_texts.append(slide_text.strip())
    
    return slide_texts

# スライド間の類似度を計算
def calculate_similarity(slides_texts):
    # スライドのテキストをベクトル化
    slide_vectors = model.encode(slides_texts)

    # コサイン類似度を計算
    similarity_matrix = cosine_similarity(slide_vectors)
    return similarity_matrix

# スライド類似度を表示
def display_similarity_matrix(similarity_matrix):
    print("スライド間の類似度行列:")
    for row in similarity_matrix:
        print([f"{sim:.2f}" for sim in row])

# 使用例
if __name__ == "__main__":
    pptx_path = sys.argv[1]  # コマンドライン引数からPowerPointファイルのパスを取得
    
    # PowerPointファイルからスライドごとのテキストを抽出
    slide_texts = extract_slide_texts(pptx_path)
    
    # スライド間の類似度を計算
    similarity_matrix = calculate_similarity(slide_texts)
    
    # 類似度行列を表示
    display_similarity_matrix(similarity_matrix)

